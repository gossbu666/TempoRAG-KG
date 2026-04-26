"""Build the 16-slide final presentation for TempoRAG-KG (April 27 submission).

Targets the 30-pt rubric × 7 TA feedback items in one deck:
  1  Title                     - intro
  2  Hook (failure example)    - intro motivation
  3  RQ / IV / DV / H          - rubric Intro 5pts
  4  Related Work + Gap        - rubric RW 4pts
  5  Method - 4-condition      - rubric Method 3pts
  6  Data + QA gen logic       - TA #3
  7  KG visualization          - TA #1(i)
  8  Link-jumping mechanism    - TA #1(ii) - uses fig_link_jumping.png
  9  Headline results          - rubric Result 3pts - uses fig_2x2_ablation.png
 10  By-hop / by-scope         - rubric Result - uses fig_by_hop.png
 11  Qualitative deep-dive     - TA #4
 12  Error taxonomy            - TA #2 + #5 - uses fig_taxonomy_by_cond.png
 13  Cost non-KG vs KG         - TA #6
 14  Lessons + Future Work     - rubric Discussion 3pts + TA #7
 15  Live demo placeholder     - rubric Demo 3pts
 16  Conclusion + thanks

Run:  python3 scripts/build_final_slides.py
Output: docs/final_slides.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "docs" / "final_slides.pptx"
FIG = REPO / "docs" / "figures"

MODEL_ORDER = [
    "gpt-4.1-nano",
    "gpt-4o-mini",
    "gpt-4o",
    "llama-70b",
    "llama-8b",
    "gpt-oss-120b",
    "gpt-oss-20b",
]

# Palette - Warm academic (maroon + cream)
NAVY = RGBColor(0x7B, 0x1F, 0x2A)         # primary: deep maroon
GOLD = RGBColor(0xA8, 0x7E, 0x40)         # accent: warm gold
INK = RGBColor(0x2A, 0x18, 0x10)          # warm dark text
MUTED = RGBColor(0x8B, 0x7E, 0x72)        # warm taupe gray
GREEN = RGBColor(0x5E, 0x73, 0x40)        # olive (positive)
RED = RGBColor(0xA0, 0x48, 0x48)          # brick (negative)
GREEN_FILL = RGBColor(0xDD, 0xE2, 0xC8)   # cream-olive
BLUE_FILL = RGBColor(0xD9, 0xDD, 0xE0)    # dusty blue-gray
YELLOW_FILL = RGBColor(0xF0, 0xE8, 0xCC)  # warm cream
ORANGE_FILL = RGBColor(0xE8, 0xCD, 0xB0)  # terracotta
RED_FILL = RGBColor(0xE8, 0xD0, 0xCE)     # dusty rose
GRAY_FILL = RGBColor(0xEE, 0xE8, 0xDC)    # warm taupe fill
PAGE_BG = RGBColor(0xFA, 0xF6, 0xEE)      # cream paper background
SIDEBAR_BG = RGBColor(0xF2, 0xEA, 0xDA)   # slightly darker cream sidebar
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ────────────────────────────────────────────────────────────────────
# Primitives
# ────────────────────────────────────────────────────────────────────

def _load_summary(sweep: str, model: str) -> dict | None:
    p = REPO / "data" / "eval" / sweep / model / "summary.json"
    return json.loads(p.read_text()) if p.exists() else None


def _new_slide(prs, *, page_num: int | None = None, total: int | None = None):
    """Add a slide with cream background + footer page number."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-bleed cream background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0), Inches(0),
                            prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAGE_BG
    bg.line.fill.background()
    # Top maroon bar (thin, full-bleed)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0), Inches(0),
                             prs.slide_width, Inches(0.12))
    top.fill.solid(); top.fill.fore_color.rgb = NAVY
    top.line.fill.background()
    # Footer page number + project tag
    if page_num is not None:
        fbox = s.shapes.add_textbox(Inches(0.4), Inches(7.15),
                                    Inches(12.5), Inches(0.3))
        p = fbox.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"TempoRAG-KG  ·  AT82.05  ·  AIT 2026"
        r.font.size = Pt(9); r.font.color.rgb = MUTED
        r.font.italic = True
        # Right-aligned page num
        pbox = s.shapes.add_textbox(Inches(12.5), Inches(7.15),
                                    Inches(0.7), Inches(0.3))
        p = pbox.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{page_num} / {total}" if total else f"{page_num}"
        r.font.size = Pt(9); r.font.color.rgb = MUTED
    return s


def _title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.32),
                                   Inches(12.5), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(26)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.runs[0].font.size = Pt(13)
        p2.runs[0].font.italic = True
        p2.runs[0].font.color.rgb = MUTED
    # Gold accent bar under title
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.4), Inches(1.28),
                                 Inches(1.2), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()


def _add_table(slide, left, top, width, height, data, *, header=True,
               col_widths=None, cell_colors=None, cell_bold=None,
               cell_italic=None, cell_align=None, font_size=11):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = INK
                    if header and r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    if cell_bold and cell_bold[r][c]:
                        run.font.bold = True
                    if cell_italic and cell_italic[r][c]:
                        run.font.italic = True
                        run.font.color.rgb = MUTED
                if cell_align and cell_align[r][c]:
                    p.alignment = cell_align[r][c]
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif cell_colors and cell_colors[r][c] is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_colors[r][c]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl


def _add_bullet_box(slide, left, top, width, height, bullets,
                    font_size=13, bullet_color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (label, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        if label:
            run_label = p.add_run()
            run_label.text = f"{label}  "
            run_label.font.size = Pt(font_size)
            run_label.font.bold = True
            run_label.font.color.rgb = NAVY
        run_body = p.add_run()
        run_body.text = text
        run_body.font.size = Pt(font_size)
        run_body.font.color.rgb = bullet_color


def _add_picture(slide, path: Path, left, top, *, width=None, height=None):
    if not path.exists():
        # Placeholder rectangle
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top,
            width or Inches(8), height or Inches(4)
        )
        ph.fill.solid()
        ph.fill.fore_color.rgb = GRAY_FILL
        ph.line.color.rgb = MUTED
        tf = ph.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[ Missing: {path.name} ]"
        r.font.size = Pt(14)
        r.font.italic = True
        r.font.color.rgb = MUTED
        return
    kwargs = {}
    if width is not None:
        kwargs["width"] = width
    if height is not None:
        kwargs["height"] = height
    slide.shapes.add_picture(str(path), left, top, **kwargs)


def _footer(slide, text: str):
    fbox = slide.shapes.add_textbox(Inches(0.4), Inches(7.05),
                                    Inches(12.5), Inches(0.4))
    p = fbox.text_frame.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = MUTED


# ────────────────────────────────────────────────────────────────────
# Slides 1-2: Title + Hook
# ────────────────────────────────────────────────────────────────────

def slide_1_title(prs, page_num=0, total=17):
    # Custom title slide: full maroon left band + cream right
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Cream full background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0), Inches(0),
                            prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAGE_BG
    bg.line.fill.background()
    # Left maroon band (1/3 of slide)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0), Inches(0),
                              Inches(4.5), prs.slide_height)
    band.fill.solid(); band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    # Gold accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4.5), Inches(0),
                               Inches(0.08), prs.slide_height)
    strip.fill.solid(); strip.fill.fore_color.rgb = GOLD
    strip.line.fill.background()

    # Vertical text on left band: course code
    cbox = s.shapes.add_textbox(Inches(0.4), Inches(0.6),
                                Inches(3.7), Inches(0.5))
    p = cbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "AT82.05  ·  NLU  ·  Spring 2026"
    r.font.size = Pt(12); r.font.color.rgb = WHITE
    r.font.bold = True

    # AIT brand on left band bottom
    abox = s.shapes.add_textbox(Inches(0.4), Inches(6.4),
                                Inches(3.7), Inches(0.7))
    tf = abox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Asian Institute"
    r.font.size = Pt(13); r.font.color.rgb = WHITE
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "of Technology"
    r.font.size = Pt(13); r.font.color.rgb = WHITE
    r.font.bold = True

    # Big title (right side, top)
    box = s.shapes.add_textbox(Inches(5.0), Inches(1.6),
                               Inches(8.0), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "TempoRAG-KG"
    r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = NAVY

    # Subtitle
    sub = s.shapes.add_textbox(Inches(5.0), Inches(3.0),
                               Inches(8.0), Inches(1.5))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Temporal-Aware Knowledge-Graph"
    r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = INK
    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = "Augmented RAG for Multi-Hop QA on SEC 10-K Filings"
    r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = INK

    # Divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(5.0), Inches(4.7),
                             Inches(2.0), Inches(0.04))
    div.fill.solid(); div.fill.fore_color.rgb = GOLD
    div.line.fill.background()

    # Team list (right side, bottom)
    team = s.shapes.add_textbox(Inches(5.0), Inches(4.95),
                                Inches(8.0), Inches(2.0))
    tf = team.text_frame
    tf.word_wrap = True
    members = [
        ("Supanut Kompayak", "st126055"),
        ("Aphisit Jaemyaem", "st126130"),
        ("Dechathon Niamsa-ard", "st126235"),
        ("Kaung Hein Htet", "st126477"),
    ]
    for i, (name, sid) in enumerate(members):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = name
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = INK
        r2 = p.add_run()
        r2.text = f"   {sid}"
        r2.font.size = Pt(11); r2.font.color.rgb = MUTED

    # Date footer (right side, bottom)
    date = s.shapes.add_textbox(Inches(5.0), Inches(6.85),
                                Inches(8.0), Inches(0.4))
    p = date.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "April 27, 2026"
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = MUTED


def slide_2_hook(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "The problem in one query")

    # Question quote box
    qbox = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.1))
    qbox.fill.solid(); qbox.fill.fore_color.rgb = YELLOW_FILL
    qbox.line.color.rgb = MUTED
    tf = qbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '"Which company had higher data-center revenue in FY2024 — NVIDIA or Intel?"'
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = INK

    # Two-column comparison
    data = [
        ["Vanilla RAG (top-k cosine)", "What we want"],
        ["Retrieves NVIDIA chunks · misses Intel\n→ \"not provided in the excerpts\"  ❌",
         "Year-aware retrieval pulls Intel FY2024\n→ \"$12,817 million\"  ✅"],
    ]
    cell_colors = [[NAVY, NAVY], [RED_FILL, GREEN_FILL]]
    cell_bold = [[True, True], [False, False]]
    _add_table(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(2.5),
               data, header=True, col_widths=[Inches(6.15), Inches(6.15)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=14)

    _add_bullet_box(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4), [
        ("Two cross-cutting needs:",
         "(1) the right year's filing,  (2) the right entity (Intel) "
         "even when the question seeds with another (NVIDIA)."),
        ("Our claim",
         "—  Both temporal filtering and KG-based entity expansion are needed; "
         "neither alone closes the gap."),
    ], font_size=13)


# ────────────────────────────────────────────────────────────────────
# Slide 3: RQ / IV / DV / H
# ────────────────────────────────────────────────────────────────────

def slide_3_rq(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Research questions, hypotheses, variables")

    # RQ table
    rq = [
        ["#", "Question", "Hypothesis"],
        ["RQ1", "Does temporal filtering lift QA over vanilla RAG?",
         "H1: Year mask raises coverage on year-anchored Qs."],
        ["RQ2", "Does KG entity expansion lift QA over vanilla RAG?",
         "H2: Graph walk recovers entities cosine misses."],
        ["RQ3 ★", "Does combining temporal + KG outperform either alone?",
         "H3: Effects are additive on multi-hop queries (hop ≥ 2)."],
        ["RQ4", "Where do the lifts come from — retrieval or generation?",
         "H4: Coverage moves; F1@answered does not."],
    ]
    cell_colors = [[None]*3 for _ in range(len(rq))]
    cell_bold = [[False]*3 for _ in range(len(rq))]
    cell_colors[3][0] = YELLOW_FILL
    cell_colors[3][1] = YELLOW_FILL
    cell_colors[3][2] = YELLOW_FILL
    for c in range(3):
        cell_bold[3][c] = True
    _add_table(s, Inches(0.4), Inches(1.3), Inches(12.6), Inches(3.3),
               rq, col_widths=[Inches(0.9), Inches(5.5), Inches(6.2)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=12)

    # IV / DV
    var = [
        ["Independent variables",
         "Retrieval condition (4: L0/L1/L2/L3)  ·  LLM (7 models)  ·  hop count (1/2/3)  ·  scope"],
        ["Dependent variables",
         "Token-F1 (primary)  ·  Coverage  ·  F1@answered  ·  95% bootstrap CI (n=1000)"],
    ]
    cc = [[GRAY_FILL, None], [GRAY_FILL, None]]
    cb = [[True, False], [True, False]]
    _add_table(s, Inches(0.4), Inches(4.8), Inches(12.6), Inches(1.6),
               var, header=False,
               col_widths=[Inches(2.6), Inches(10.0)],
               cell_colors=cc, cell_bold=cb, font_size=12)

    _footer(s, "129 questions · 3,612 cached predictions · single-pass evaluation per cell.")


# ────────────────────────────────────────────────────────────────────
# Slide 4: Related Work + Gap
# ────────────────────────────────────────────────────────────────────

def slide_4_related_work(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Related work — and the gap we fill",
           "Two retrieval axes (entity-aware × time-aware)")

    matrix = [
        ["", "Time-blind", "Time-aware"],
        ["Entity-blind",
         "Vanilla RAG (cosine top-k)\nLewis et al. 2020",
         "TA-RAG (year mask only)\n— our L1 baseline"],
        ["Entity-aware",
         "KG²RAG (entity walk only)\nZhu et al. 2025  — our L2",
         "TempoRAG-KG  ★\nThis work  — our L3"],
    ]
    cell_colors = [[None]*3 for _ in range(3)]
    cell_bold = [[False]*3 for _ in range(3)]
    for r in range(3):
        cell_colors[r][0] = GRAY_FILL
        cell_bold[r][0] = True
    for c in range(3):
        cell_colors[0][c] = GRAY_FILL
        cell_bold[0][c] = True
    cell_colors[2][2] = YELLOW_FILL
    cell_bold[2][2] = True
    _add_table(s, Inches(1.5), Inches(1.5), Inches(10.0), Inches(3.3),
               matrix, header=False,
               col_widths=[Inches(2.5), Inches(3.75), Inches(3.75)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=14)

    _add_bullet_box(s, Inches(0.4), Inches(5.0), Inches(12.6), Inches(2.0), [
        ("Gap",
         "No prior work tests entity × time as an explicit 2 × 2 factorial. "
         "KG²RAG ignores time; TA-RAG ignores entities."),
        ("Our contribution",
         "Decompose the question into a 2 × 2 ablation; isolate the marginal "
         "value of each axis on multi-hop temporal QA."),
        ("Domain twist",
         "10-K filings are date-stamped by construction — a clean substrate "
         "to test temporal mechanisms (vs. Wikipedia where dates are noisy)."),
    ], font_size=13)


# ────────────────────────────────────────────────────────────────────
# Slide 5: Method — 4-condition pipeline
# ────────────────────────────────────────────────────────────────────

def slide_5_method(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Method — 4-condition retrieval ablation",
           "Same chunks · same prompts · same models · only retrieval differs")

    cond = [
        ["Cond", "Time mask", "KG walk", "What it tests"],
        ["L0  Vanilla", "—", "—", "Pure cosine top-k baseline"],
        ["L1  TimeFilter", "✅ year[]", "—", "Effect of temporal filtering only"],
        ["L2  KG²RAG", "—", "✅ entity expand", "Effect of entity expansion only"],
        ["L3  TempoRAG-KG", "✅ year[]", "✅ entity expand", "Combined effect — additivity test"],
    ]
    cell_colors = [[None]*4 for _ in range(5)]
    cell_bold = [[False]*4 for _ in range(5)]
    cell_colors[4] = [YELLOW_FILL]*4
    for c in range(4):
        cell_bold[4][c] = True
    _add_table(s, Inches(0.4), Inches(1.5), Inches(12.6), Inches(3.0),
               cond,
               col_widths=[Inches(2.4), Inches(2.0), Inches(2.6),
                           Inches(5.6)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=13)

    _add_bullet_box(s, Inches(0.4), Inches(4.7), Inches(12.6), Inches(2.4), [
        ("Pipeline",
         "Question → embed → cosine top-N → (optional year mask) → "
         "(optional KG entity-walk expansion) → top-k=5 chunks → LLM answer."),
        ("Why 5 chunks",
         "Pre-piloted; 5 fits all 7 models' context budgets and matches "
         "KG²RAG's reported optimum."),
        ("Determinism",
         "All retrieval is cached on (question_id, condition); identical "
         "inputs ⇒ identical outputs."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 6: Data + QA generation logic (TA #3)
# ────────────────────────────────────────────────────────────────────

def slide_6_data(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Data — corpus & question construction",
           "TA #3: how were the 129 questions built?")

    # Corpus + KG stats
    corpus = [
        ["Corpus", "10 issuers × FY2019–2024  →  7,467 chunks"],
        ["Source", "SEC EDGAR — 10-K & 10-Q filings"],
        ["Chunking", "1500-char windows, 200-char overlap"],
        ["KG",  "57,718 triples · 293 unique subjects · 26k predicates"],
        ["KG extractor", "gpt-4.1-nano · temporal-aware schema (subj, pred, obj, valid_from, valid_to)"],
    ]
    cc = [[GRAY_FILL, None] for _ in range(len(corpus))]
    cb = [[True, False] for _ in range(len(corpus))]
    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(2.6),
               corpus, header=False,
               col_widths=[Inches(2.4), Inches(10.2)],
               cell_colors=cc, cell_bold=cb, font_size=12)

    _add_bullet_box(s, Inches(0.4), Inches(4.2), Inches(12.6), Inches(2.9), [
        ("Question source",
         "FinReflectKG-MultiHop benchmark, filtered to our issuer × "
         "year coverage  →  129 questions."),
        ("Decomposition",
         "1-hop (single chunk)  ·  2-hop (cross-chunk same year)  ·  "
         "3-hop (cross-year or cross-company)."),
        ("Scope tags",
         "intra · inter_year · cross_company · fiscal_vs_calendar · forward_looking."),
        ("Synthetic pool",
         "128 LLM-generated candidates → 4-axis auto-vet → only 4 passed "
         "(3.1%); dropped from evaluation, kept as a methodological lesson."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 7: KG visualization (TA #1i)
# ────────────────────────────────────────────────────────────────────

def slide_7_kg_viz(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "What does the KG look like?",
           "TA #1(i): structure & sample subgraph")

    # Stats grid
    stats = [
        ["Triples", "Unique subjects", "Filings covered", "Avg triples/chunk"],
        ["57,718", "293", "60", "7.7"],
    ]
    cell_colors = [[NAVY]*4, [BLUE_FILL]*4]
    cell_bold = [[True]*4, [True]*4]
    cell_align = [[PP_ALIGN.CENTER]*4, [PP_ALIGN.CENTER]*4]
    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(1.4),
               stats, header=True,
               col_widths=[Inches(3.15)]*4,
               cell_colors=cell_colors, cell_bold=cell_bold,
               cell_align=cell_align, font_size=18)

    # Sample subgraph (text-based illustration)
    sub = s.shapes.add_textbox(Inches(0.4), Inches(3.0),
                               Inches(12.6), Inches(0.4))
    p = sub.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Sample subgraph for AAPL FY2022 (4 of ~1,200 triples):"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY

    sample = [
        ["Subject", "Predicate", "Object", "valid_from", "valid_to"],
        ["Apple Inc.", "reported_revenue", "$394,328 million", "2021-09-26", "2022-09-24"],
        ["iPhone segment", "contributed_revenue_of", "$205,489 million", "2021-09-26", "2022-09-24"],
        ["Tim Cook", "is_ceo_of", "Apple Inc.", "2011-08-24", None],
        ["Apple Inc.", "headquartered_in", "Cupertino, California", None, None],
    ]
    cell_colors = [[None]*5 for _ in range(len(sample))]
    cell_italic = [[False]*5 for _ in range(len(sample))]
    for r in range(1, len(sample)):
        for c in range(3, 5):
            cell_italic[r][c] = True
    _add_table(s, Inches(0.4), Inches(3.5), Inches(12.6), Inches(2.6),
               sample, col_widths=[Inches(2.6), Inches(2.8),
                                   Inches(3.6), Inches(1.8), Inches(1.8)],
               cell_colors=cell_colors, cell_italic=cell_italic, font_size=11)

    _footer(s, "valid_from / valid_to are extracted explicitly per triple — used by L3 hard mask.")


# ────────────────────────────────────────────────────────────────────
# Slide 8: Link-jumping mechanism (TA #1ii)
# ────────────────────────────────────────────────────────────────────

def slide_8_link_jumping(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "How does the KG \"jump\" between chunks?",
           "TA #1(ii): link-jumping mechanism — seed → entity → expanded chunks")

    # Use the figure if available
    _add_picture(s, FIG / "fig_link_jumping.png",
                 Inches(0.4), Inches(1.4),
                 width=Inches(7.8))

    _add_bullet_box(s, Inches(8.4), Inches(1.4), Inches(4.8), Inches(5.5), [
        ("Step 1 — Seed",
         "Cosine top-3 chunks from the question."),
        ("Step 2 — Entity link",
         "Extract subjects/objects from each seed's KG triples."),
        ("Step 3 — Expand",
         "Pull every other chunk that mentions any seed entity."),
        ("Step 4 — Re-rank",
         "Cosine on the expanded pool; keep top-5."),
        ("L3 only — Time mask",
         "Drop expanded chunks whose triples' [valid_from, valid_to] "
         "don't overlap the query's year[] filter."),
    ], font_size=13)


# ────────────────────────────────────────────────────────────────────
# Slide 9: Headline results
# ────────────────────────────────────────────────────────────────────

def slide_9_headline(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Headline results — three findings",
           "7 models × 129 questions × 4 conditions = 3,612 cached predictions")

    _add_picture(s, FIG / "fig_2x2_ablation.png",
                 Inches(0.4), Inches(1.4),
                 width=Inches(7.8))

    _add_bullet_box(s, Inches(8.4), Inches(1.4), Inches(4.8), Inches(5.5), [
        ("① L1 universal lift",
         "TimeFilter beats Vanilla on 7/7 models  (+5.6% to +24.7%, "
         "avg +15.9%).  H1 supported."),
        ("② L2 universal regression",
         "KG²RAG-only loses on 7/7 models  (avg −6.7%).  H2 refuted — "
         "graph walk without time injects noise on this corpus."),
        ("③ L3 hop-3 sweet spot",
         "TempoRAG-KG averages ≈ L1, but on hop=3 / gpt-4.1-nano: "
         "+0.071 over L1.  H3 partially supported."),
        ("④ Bottleneck = generation",
         "F1@answered ≈ flat at 0.36; coverage moves with retrieval. "
         "H4 supported."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 9b: Ablation table — concrete numbers for all 7 × 4 cells
# ────────────────────────────────────────────────────────────────────

def slide_9b_ablation_table(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Ablation table — 7 models × 4 conditions",
           "Token-F1 mean · best per row highlighted · Δ% over Vanilla baseline")

    rows = [["Model", "L0 Vanilla", "L1 TimeFilter", "L2 KG²RAG",
             "L3 TempoRAG-KG", "Δ%(L1−L0)", "Δ%(L3−L1)"]]
    cell_colors = [[None]*7]
    cell_bold = [[False]*7]
    cell_italic = [[False]*7]

    sums = {"vanilla": 0, "timefilter": 0, "kg2rag": 0, "temporag": 0}
    n = 0

    for m in MODEL_ORDER:
        f = {c: _load_summary(c, m) for c in
             ("vanilla", "timefilter", "kg2rag", "temporag")}
        if any(v is None for v in f.values()):
            continue
        l0 = f["vanilla"]["f1_mean"]
        l1 = f["timefilter"]["f1_mean"]
        l2 = f["kg2rag"]["f1_mean"]
        l3 = f["temporag"]["f1_mean"]
        d1 = (l1 - l0) / l0 * 100 if l0 else 0
        d3 = (l3 - l1) / l1 * 100 if l1 else 0
        sums["vanilla"] += l0; sums["timefilter"] += l1
        sums["kg2rag"] += l2; sums["temporag"] += l3
        n += 1

        vals = [l0, l1, l2, l3]
        best_idx = vals.index(max(vals))  # 0..3 → cols 1..4
        rows.append([m, f"{l0:.3f}", f"{l1:.3f}", f"{l2:.3f}", f"{l3:.3f}",
                     f"{d1:+.1f}%", f"{d3:+.1f}%"])
        rc = [None]*7
        rb = [False]*7
        ri = [False]*7
        rc[1 + best_idx] = GREEN_FILL
        rb[1 + best_idx] = True
        rc[5] = GREEN_FILL if d1 > 0 else RED_FILL
        rc[6] = GREEN_FILL if d3 > 0 else RED_FILL
        cell_colors.append(rc)
        cell_bold.append(rb)
        cell_italic.append(ri)

    if n:
        avg = {k: v / n for k, v in sums.items()}
        d1 = (avg["timefilter"] - avg["vanilla"]) / avg["vanilla"] * 100
        d3 = (avg["temporag"] - avg["timefilter"]) / avg["timefilter"] * 100
        vals = list(avg.values())
        best_idx = vals.index(max(vals))
        rows.append(["AVG (n=7)",
                     f"{avg['vanilla']:.3f}", f"{avg['timefilter']:.3f}",
                     f"{avg['kg2rag']:.3f}", f"{avg['temporag']:.3f}",
                     f"{d1:+.1f}%", f"{d3:+.1f}%"])
        rc = [GRAY_FILL]*7
        rb = [True]*7
        rc[1 + best_idx] = ORANGE_FILL
        cell_colors.append(rc)
        cell_bold.append(rb)
        cell_italic.append([False]*7)

    align = [[None]*7 for _ in range(len(rows))]
    for r in range(len(rows)):
        for c in range(1, 7):
            align[r][c] = PP_ALIGN.RIGHT

    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(3.5),
               rows,
               col_widths=[Inches(2.0), Inches(1.5), Inches(1.7),
                           Inches(1.5), Inches(1.9), Inches(1.5),
                           Inches(1.5)],
               cell_colors=cell_colors, cell_bold=cell_bold,
               cell_italic=cell_italic, cell_align=align, font_size=11)

    _add_bullet_box(s, Inches(0.4), Inches(5.1), Inches(12.6), Inches(2.0), [
        ("L1 wins 7/7 models",
         "TimeFilter is the best column on every row — universal lift."),
        ("L2 loses 7/7 models",
         "KG²RAG-only regresses below L0 in every cell — graph noise "
         "without temporal anchor."),
        ("L3 ≈ L1 overall, wins on hop = 3",
         "The averaged story masks L3's targeted win on hop=3 / "
         "gpt-4.1-nano (+0.071 over L1) — see next slide."),
        ("Analyses performed",
         "① overall ablation  ② by-hop  ③ by-scope  ④ failure taxonomy "
         "⑤ qualitative case  ⑥ cost — covered in slides 10-13."),
    ], font_size=11)


# ────────────────────────────────────────────────────────────────────
# Slide 10: By hop / by scope
# ────────────────────────────────────────────────────────────────────

def slide_10_by_hop(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Where L3 wins — multi-hop and cross-company",
           "By-hop and by-scope breakdown")

    _add_picture(s, FIG / "fig_by_hop.png",
                 Inches(0.4), Inches(1.4),
                 width=Inches(6.4))
    _add_picture(s, FIG / "fig_by_scope.png",
                 Inches(7.0), Inches(1.4),
                 width=Inches(6.0))

    _add_bullet_box(s, Inches(0.4), Inches(5.5), Inches(12.6), Inches(1.6), [
        ("hop = 3",
         "L3 +0.071 over L1 on gpt-4.1-nano  →  exactly where multi-hop "
         "entity bridging is supposed to help."),
        ("hop = 1",
         "All conditions converge — vanilla cosine already finds the "
         "single chunk; year mask and KG add nothing."),
        ("Scope cross_company",
         "L1 alone gives the headline lift (year mask pulls the "
         "second company's filing into top-k)."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 11: Qualitative deep-dive (TA #4)
# ────────────────────────────────────────────────────────────────────

def slide_11_qualitative(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Qualitative deep-dive — one question, four conditions",
           "TA #4: NVDA vs Intel data-center revenue, FY2024  ·  model = gpt-4o-mini")

    walk = [
        ["Cond", "Top-5 chunks include Intel?", "Answer", "Mechanism"],
        ["L0  Vanilla", "❌  NVDA-heavy",
         "\"not provided\"  ❌",
         "Cosine pulls semantically-similar NVDA chunks; INTC misses top-5."],
        ["L1  TimeFilter", "✅  FY2024 mask lifts INTC",
         "\"$12,817M\"  ✅",
         "Year mask forces FY2024 filings into the candidate pool."],
        ["L2  KG²RAG", "❌  identical to L0",
         "\"not provided\"  ❌",
         "KG entity walk seeds with NVDA, expands NVDA-related entities — never bridges to INTC."],
        ["L3  TempoRAG-KG", "✅  identical to L1",
         "\"$12,817M\"  ✅",
         "Time mask carries the lift; KG adds nothing on this query."],
    ]
    cell_colors = [[None]*4 for _ in range(5)]
    cell_bold = [[False]*4 for _ in range(5)]
    cell_colors[1] = [None, RED_FILL, RED_FILL, None]
    cell_colors[2] = [None, GREEN_FILL, GREEN_FILL, None]
    cell_colors[3] = [None, RED_FILL, RED_FILL, None]
    cell_colors[4] = [None, GREEN_FILL, GREEN_FILL, None]
    for r in range(1, 5):
        cell_bold[r][0] = True
    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(4.4),
               walk, col_widths=[Inches(1.9), Inches(2.8),
                                 Inches(2.0), Inches(5.9)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=11)

    _add_bullet_box(s, Inches(0.4), Inches(5.9), Inches(12.6), Inches(1.2), [
        ("Cache hits prove identity",
         "L2's prompt is byte-identical to L0's (same retrieved chunks); "
         "L3's matches L1's. Temporal is the active ingredient on this query."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 12: Error taxonomy (TA #2 + #5)
# ────────────────────────────────────────────────────────────────────

def slide_12_taxonomy(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Failure taxonomy — where do the errors come from?",
           "TA #2 + #5: model-level vs corpus-level failure modes")

    _add_picture(s, FIG / "fig_taxonomy_by_cond.png",
                 Inches(0.4), Inches(1.3),
                 width=Inches(7.4))

    _add_bullet_box(s, Inches(8.0), Inches(1.3), Inches(5.2), Inches(5.5), [
        ("A4  IDK-when-answerable",
         "41.8% of all 3,607 classified predictions  →  the dominant "
         "failure mode across all conditions."),
        ("A1  Tersification artefact",
         "Token-F1 floor: gold has 8 tokens, model gives 2-3 correct "
         "ones — F1 ≈ 0.4 even when fact is right."),
        ("A2  Wrong-year retrieval",
         "Drops from 18% (L0) to 6% (L1)  ←  TimeFilter directly fixes this."),
        ("B1-B5  Corpus-level limits",
         "10-K cannot answer stock prices, forward guidance, opinions  "
         "— hard ceiling not addressable by any retrieval."),
        ("Reliability",
         "Cohen's κ = 0.200 (LLM-vs-LLM, n=20). Slight agreement; "
         "framed honestly as not a human IRR substitute."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 13: Cost (TA #6)
# ────────────────────────────────────────────────────────────────────

def slide_13_cost(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Cost — non-KG vs KG conditions",
           "TA #6: concrete dollars per pipeline component")

    cost = [
        ["Component", "One-time", "Recurring (per Q)"],
        ["Corpus chunking + embedding (one-time)", "$0.40", "—"],
        ["KG extraction (7,467 chunks · gpt-4.1-nano)", "$8.20", "—"],
        ["L0/L1 evaluation (cached)", "—", "$0  (cache hit)"],
        ["L2/L3 evaluation (cached)", "—", "$0  (cache hit)"],
        ["Failure-taxonomy LLM judge (gpt-4o-mini, n=1,183)", "$0.30", "—"],
        ["Streamlit demo (live, ~5¢ per uncached query)", "—", "$0.05"],
    ]
    cell_colors = [[None]*3 for _ in range(len(cost))]
    cell_bold = [[False]*3 for _ in range(len(cost))]
    # Highlight KG vs non-KG
    cell_colors[2][0] = ORANGE_FILL  # KG extract row
    cell_colors[2][1] = ORANGE_FILL
    cell_bold[2][1] = True
    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.6), Inches(3.4),
               cost, col_widths=[Inches(7.2), Inches(2.7), Inches(2.7)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=12)

    _add_bullet_box(s, Inches(0.4), Inches(5.0), Inches(12.6), Inches(2.0), [
        ("Total project cost",
         "$8.90  (one-time)  +  ~$0.05 per uncached demo query."),
        ("KG premium",
         "$8.20 = 92% of total. KG buys L2/L3 capability at a 21× "
         "multiple over L0/L1 setup cost."),
        ("Per-question cost in production",
         "$0  on cached queries, $0.05 on uncached — well under any "
         "operational threshold."),
    ], font_size=13)


# ────────────────────────────────────────────────────────────────────
# Slide 14: Lessons + Future Work (TA #7)
# ────────────────────────────────────────────────────────────────────

def slide_14_lessons(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "What we learned · Future work",
           "Discussion 3pts + TA #7 inject-large/retrieve-small")

    _add_bullet_box(s, Inches(0.4), Inches(1.3), Inches(12.6), Inches(2.8), [
        ("Lesson 1 — Negative findings matter",
         "L2 KG-only regresses universally. Most papers hide this; we "
         "report it because the mechanism (KG-induced noise without "
         "temporal anchor) is the actual contribution."),
        ("Lesson 2 — Token-F1 has a ceiling",
         "0/129 questions reach F1 ≥ 0.8; tersification dominates. "
         "Future work: scoring against entity-set + canonical numeric value."),
        ("Lesson 3 — Synth QA is hard",
         "Strict 4-axis auto-vet rejected 124/128 candidates. "
         "Lesson: hop-count is verifiable; relevance + answerability are not, "
         "without human labels."),
    ], font_size=12)

    # Future work box
    fbox = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(4.4), Inches(12.6), Inches(2.5))
    fbox.fill.solid(); fbox.fill.fore_color.rgb = BLUE_FILL
    fbox.line.color.rgb = NAVY

    tbox = s.shapes.add_textbox(Inches(0.6), Inches(4.5),
                                Inches(12.2), Inches(0.4))
    p = tbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Future work — TA #7 inject-large / retrieve-small"
    r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = NAVY

    _add_bullet_box(s, Inches(0.6), Inches(4.95), Inches(12.2), Inches(2.0), [
        ("Two-tier model split",
         "Large LLM at extraction (one-time, $$); small LLM at retrieval-"
         "and-answer (per-query, $).  Predicted Pareto improvement."),
        ("Generation-side intervention",
         "A4 IDK 41.8% says retrieval is solved; the next bottleneck is "
         "uncertainty-aware decoding."),
    ], font_size=12)


# ────────────────────────────────────────────────────────────────────
# Slide 15: Live demo placeholder
# ────────────────────────────────────────────────────────────────────

def slide_15_demo(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "🎬  Live demo — Streamlit")

    # Big centered placeholder for screen recording / screenshot
    ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0.6), Inches(1.4),
                            Inches(8.5), Inches(5.4))
    ph.fill.solid(); ph.fill.fore_color.rgb = GRAY_FILL
    ph.line.color.rgb = MUTED
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[ Switch to live Streamlit window ]"
    r.font.size = Pt(20); r.font.italic = True; r.font.color.rgb = MUTED
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    p2.text = "streamlit run app/streamlit_app.py    (port 8501)"
    p2.runs[0].font.size = Pt(13); p2.runs[0].font.color.rgb = MUTED

    # Demo plan on the right
    dbox = s.shapes.add_textbox(Inches(9.4), Inches(1.4),
                                Inches(3.8), Inches(0.4))
    p = dbox.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Demo plan — 60 sec"
    r.font.bold = True; r.font.size = Pt(15); r.font.color.rgb = NAVY

    _add_bullet_box(s, Inches(9.4), Inches(1.85), Inches(3.8), Inches(5.0), [
        ("0:00", "Open app → click sample “NVDA vs Intel data-centre”"),
        ("0:10", "Pick year [2024] → Compare all 4 conditions → Ask"),
        ("0:25", "L0 ❌  · L1 ✅  · L2 ❌  · L3 ✅ + 🧬 KG-expanded badge"),
        ("0:45", "Show retrieved chunks: Intel only in L1/L3"),
        ("0:55", "Wrap: temporal carries the lift; KG visible on L3"),
    ], font_size=11)


# ────────────────────────────────────────────────────────────────────
# Slide 16: Conclusion
# ────────────────────────────────────────────────────────────────────

def slide_16_conclusion(prs, page_num=0, total=17):
    s = _new_slide(prs, page_num=page_num, total=total)
    _title(s, "Conclusion")

    _add_bullet_box(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(4.5), [
        ("①  Temporal filtering is universally beneficial",
         "+15.9% avg F1, all 7 models, no regressions. H1 supported."),
        ("②  KG entity expansion alone regresses",
         "−6.7% avg F1. H2 refuted. Negative finding worth reporting."),
        ("③  Combined L3 wins where it should — hop = 3",
         "+0.071 over L1 on gpt-4.1-nano. H3 supported on the targeted slice."),
        ("④  The next bottleneck is generation, not retrieval",
         "A4 IDK = 41.8%; F1@answered flat. H4 supported."),
        ("Open-sourced",
         "Code, KG, predictions, and Streamlit app: github.com/gossbu666/TempoRAG-KG"),
    ], font_size=14)

    # Thank-you line
    thanks = s.shapes.add_textbox(Inches(0.4), Inches(6.4),
                                  Inches(12.5), Inches(0.5))
    p = thanks.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thank you  ·  Questions?"
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = NAVY


# ────────────────────────────────────────────────────────────────────
# Main
# ────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_1_title,
        slide_2_hook,
        slide_3_rq,
        slide_4_related_work,
        slide_5_method,
        slide_6_data,
        slide_7_kg_viz,
        slide_8_link_jumping,
        slide_9_headline,
        slide_9b_ablation_table,
        slide_10_by_hop,
        slide_11_qualitative,
        slide_12_taxonomy,
        slide_13_cost,
        slide_14_lessons,
        slide_15_demo,
        slide_16_conclusion,
    ]
    total = len(builders)
    for i, fn in enumerate(builders, 1):
        fn(prs, page_num=i, total=total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(REPO)}  ({OUT.stat().st_size:,} bytes, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
