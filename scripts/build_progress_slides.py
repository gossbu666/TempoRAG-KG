"""Build the 7-slide progress presentation from live eval summaries.

Re-run after TimeFilter / vanilla numbers change; the PPTX regenerates from
`data/eval/{vanilla,timefilter}/<model>/summary.json`. Pipeline diagram
(slide 3) is left as a placeholder rectangle — paste a PNG exported from
`docs/progress_report/pipeline.drawio`.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "docs" / "progress_report" / "progress_report.pptx"

MODEL_ORDER = [
    "gpt-4.1-nano",
    "gpt-4o-mini",
    "gpt-4o",
    "llama-70b",
    "llama-8b",
    "gpt-oss-120b",
    "gpt-oss-20b",
]
SMALL_TIER = {"llama-8b", "gpt-oss-20b"}  # highlight for RQ4 capability story

NAVY = RGBColor(0x1F, 0x3A, 0x68)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x77, 0x77, 0x77)
GREEN_FILL = RGBColor(0xD5, 0xE8, 0xD4)
BLUE_FILL = RGBColor(0xDA, 0xE8, 0xFC)
YELLOW_FILL = RGBColor(0xFF, 0xF2, 0xCC)
ORANGE_FILL = RGBColor(0xFF, 0xE6, 0xCC)
GRAY_FILL = RGBColor(0xEE, 0xEE, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _load_summary(sweep: str, model: str) -> dict | None:
    p = REPO / "data" / "eval" / sweep / model / "summary.json"
    return json.loads(p.read_text()) if p.exists() else None


def _title(slide, text: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(26)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.runs[0].font.size = Pt(13)
        p2.runs[0].font.italic = True
        p2.runs[0].font.color.rgb = MUTED


def _add_table(slide, left, top, width, height, data, *, header=True,
               col_widths=None, cell_colors=None, cell_bold=None,
               cell_italic=None, cell_align=None, font_size=11):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = INK
                    if header and r == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                    if cell_bold and cell_bold[r][c]:
                        run.font.bold = True
                    if cell_italic and cell_italic[r][c]:
                        run.font.italic = True
                        run.font.color.rgb = MUTED
                if cell_align and cell_align[r][c]:
                    p.alignment = cell_align[r][c]
            if header and r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
            elif cell_colors and cell_colors[r][c] is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_colors[r][c]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return tbl


def _add_bullet_box(slide, left, top, width, height, bullets,
                    font_size=13, bullet_color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (label, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        run_label = p.add_run()
        run_label.text = f"{label}  "
        run_label.font.size = Pt(font_size)
        run_label.font.bold = True
        run_label.font.color.rgb = NAVY
        run_body = p.add_run()
        run_body.text = text
        run_body.font.size = Pt(font_size)
        run_body.font.color.rgb = bullet_color


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def slide_1_recap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "From Wikipedia QA to SEC 10-K filings",
           "Dataset pivot: anchoring problem → year-stamped corpus")
    data = [
        ["", "v1  (original proposal)", "v2  (this sprint)"],
        ["Dataset",
         "HotpotQA + MuSiQue\n(Wikipedia multi-hop)",
         "25 × 10-K filings\n(5 tech mega-caps × 5 fiscal years)"],
        ["Core limitation we hit",
         "Facts lack explicit year anchors;\nWikipedia spans decades with no canonical\n\"as of\" date",
         "Every 10-K is filed for a specific fiscal year\n— year is metadata, not inference"],
        ["Why it matters for us",
         "Can't evaluate temporal filtering when\nground-truth time itself is ambiguous",
         "Temporal correctness is testable —\nfilter right year, get right fact"],
    ]
    col_widths = [Inches(2.2), Inches(5.3), Inches(5.3)]
    cell_colors = [[None]*3 for _ in range(len(data))]
    cell_bold = [[False]*3 for _ in range(len(data))]
    for r in range(1, len(data)):
        cell_colors[r][0] = GRAY_FILL
        cell_bold[r][0] = True
        cell_colors[r][1] = None
        cell_colors[r][2] = BLUE_FILL
    _add_table(s, Inches(0.4), Inches(1.4), Inches(12.8), Inches(5.3),
               data, col_widths=col_widths, cell_colors=cell_colors,
               cell_bold=cell_bold, font_size=12)


def slide_2_rqs(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Research questions & experimental design")
    rq_data = [
        ["#", "Question"],
        ["RQ1", "Where does KG²RAG fail on temporal multi-hop QA?"],
        ["RQ2", "Does TempoRAG-KG lift F1 over KG²RAG?"],
        ["RQ3", "How accurate is temporal extraction?"],
        ["RQ4 ★", "Can small-LM + temporal-KG ≈ large-LM alone?"],
    ]
    cell_colors = [[None]*2 for _ in range(len(rq_data))]
    cell_bold = [[False]*2 for _ in range(len(rq_data))]
    cell_colors[4][0] = YELLOW_FILL
    cell_colors[4][1] = YELLOW_FILL
    cell_bold[4][0] = True
    cell_bold[4][1] = True
    _add_table(s, Inches(0.4), Inches(1.25), Inches(12.8), Inches(2.0),
               rq_data, col_widths=[Inches(1.2), Inches(11.6)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=13)

    # Hypothesis line
    hbox = s.shapes.add_textbox(Inches(0.4), Inches(3.45), Inches(12.8), Inches(0.6))
    tf = hbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "RQ4 hypothesis:  "
    r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = NAVY
    r = p.add_run()
    r.text = ("Temporal grounding closes the small-vs-large capability "
              "gap on temporal QA.")
    r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = INK

    # 2x2 matrix
    mbox = s.shapes.add_textbox(Inches(0.4), Inches(4.1), Inches(4.0), Inches(0.4))
    tf = mbox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Test: 2 × 2 factorial · 7 models · 129 Qs"
    r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = NAVY

    matrix = [
        ["", "No KG", "+ KG"],
        ["No Temporal", "Vanilla  ✅", "KG²RAG  ⏳"],
        ["+ Temporal", "TimeFilter  ✅", "TempoRAG-KG  ⏳"],
    ]
    cell_colors = [[None]*3 for _ in range(3)]
    cell_bold = [[False]*3 for _ in range(3)]
    cell_colors[1][1] = GREEN_FILL
    cell_colors[2][1] = GREEN_FILL
    cell_colors[1][2] = ORANGE_FILL
    cell_colors[2][2] = ORANGE_FILL
    cell_bold[0][0] = True
    cell_bold[0][1] = True
    cell_bold[0][2] = True
    cell_bold[1][0] = True
    cell_bold[2][0] = True
    cell_colors[0][0] = GRAY_FILL
    cell_colors[1][0] = GRAY_FILL
    cell_colors[2][0] = GRAY_FILL
    _add_table(s, Inches(0.4), Inches(4.55), Inches(7.8), Inches(2.0),
               matrix, header=False,
               col_widths=[Inches(2.2), Inches(2.8), Inches(2.8)],
               cell_colors=cell_colors, cell_bold=cell_bold, font_size=14)

    # Footnote
    fbox = s.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.8), Inches(0.4))
    fp = fbox.text_frame.paragraphs[0]
    fp.text = ("Today: L0 + L1 lanes complete.  "
               "KG-based cells after extract completes.")
    fp.runs[0].font.size = Pt(11)
    fp.runs[0].font.italic = True
    fp.runs[0].font.color.rgb = MUTED


def slide_3_pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Pipeline — what we built this sprint")
    # Placeholder box for pipeline.drawio export
    ph = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(0.8), Inches(1.4),
                            Inches(11.8), Inches(5.5))
    ph.fill.solid()
    ph.fill.fore_color.rgb = GRAY_FILL
    ph.line.color.rgb = MUTED
    tf = ph.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "[ Paste pipeline.drawio PNG export here ]"
    r.font.size = Pt(18)
    r.font.italic = True
    r.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "File: docs/progress_report/pipeline.drawio"
    p2.runs[0].font.size = Pt(12)
    p2.runs[0].font.color.rgb = MUTED


def _build_results_table():
    rows = [["Model", "Vanilla F1", "+TimeFilter", "Δ F1 %",
             "Van cov", "TF cov", "Δ cov", "F1@ans V→T"]]
    cell_colors = [[None]*8]
    cell_bold = [[False]*8]
    cell_italic = [[False]*8]
    for m in MODEL_ORDER:
        v = _load_summary("vanilla", m)
        t = _load_summary("timefilter", m)
        if v is None:
            continue
        vf, vc = v["f1_mean"], v.get("coverage", 0)
        vfa = v.get("f1_answered_mean", 0)
        if t is None:
            rows.append([m, f"{vf:.3f}", "—", "—",
                         f"{vc:.1%}", "—", "—", f"{vfa:.3f} → —"])
            cell_colors.append([None]*8)
            cell_bold.append([False]*8)
            cell_italic.append([False]*8)
            continue
        tf, tc = t["f1_mean"], t.get("coverage", 0)
        tfa = t.get("f1_answered_mean", 0)
        d = (tf - vf) / vf * 100 if vf else 0
        dc = (tc - vc) * 100
        rows.append([m, f"{vf:.3f}", f"{tf:.3f}", f"{d:+.1f}%",
                     f"{vc:.1%}", f"{tc:.1%}", f"{dc:+.1f}pp",
                     f"{vfa:.3f} → {tfa:.3f}"])
        row_colors = [None]*8
        row_bold = [False]*8
        row_italic = [False]*8
        # Green on Δ F1% column for all lift rows
        if d > 0:
            row_colors[3] = GREEN_FILL
            row_bold[3] = True
        # Blue-bold on Δ cov
        row_colors[6] = BLUE_FILL
        row_bold[6] = True
        # Gray-italic on F1@ans column to show it's flat
        row_italic[7] = True
        # Small-tier yellow tint on model name
        if m in SMALL_TIER:
            row_colors[0] = YELLOW_FILL
            row_bold[0] = True
        cell_colors.append(row_colors)
        cell_bold.append(row_bold)
        cell_italic.append(row_italic)
    return rows, cell_colors, cell_bold, cell_italic


def slide_5_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Temporal filtering lifts F1 on all 7 models  (+5.6% to +24.7%)",
           "L0 Vanilla vs L1 TimeFilter · 129 Qs · 95% bootstrap CIs available")
    data, colors, bolds, italics = _build_results_table()
    cell_align = [[None]*8 for _ in range(len(data))]
    for r in range(len(data)):
        for c in range(1, 8):
            cell_align[r][c] = PP_ALIGN.RIGHT
    _add_table(s, Inches(0.3), Inches(1.4), Inches(12.9), Inches(3.6),
               data, col_widths=[Inches(1.8), Inches(1.3), Inches(1.4),
                                 Inches(1.3), Inches(1.3), Inches(1.3),
                                 Inches(1.3), Inches(2.2)],
               cell_colors=colors, cell_bold=bolds, cell_italic=italics,
               cell_align=cell_align, font_size=11)

    _add_bullet_box(s, Inches(0.4), Inches(5.3), Inches(12.6), Inches(1.9), [
        ("Universal lift",
         "— 7/7 models improved (+5.6% to +24.7% F1)."),
        ("Coverage drives F1",
         "— Δ cov +3 to +10pp;  F1@answered stays flat ≈ 0.35.  "
         "Bottleneck is retrieval, not answer quality."),
        ("Small-model benefit",
         "— llama-8b +24.7%, gpt-oss-20b +21.0% vs gpt-4o-mini +7.8%.  "
         "Early signal for RQ4 capability parity."),
    ], font_size=12)


def _slice_f1(path, qa_meta):
    from collections import defaultdict
    slices = defaultdict(list)
    for line in open(path):
        r = json.loads(line)
        m = qa_meta.get(r["question_id"], {})
        slices[("hop", m.get("hop"))].append(r["f1"])
        slices[("scope", m.get("scope"))].append(r["f1"])
    return slices


def _load_qa_meta():
    meta = {}
    for p in [REPO / "data" / "qa" / "home_grown.jsonl",
              REPO / "data" / "qa" / "multihop_filtered.jsonl"]:
        for line in p.open():
            r = json.loads(line)
            meta[r["question_id"]] = {
                "hop": r.get("hop_count"),
                "scope": r.get("scope"),
            }
    return meta


def slide_6_ceiling(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Where TimeFilter hits its ceiling",
           "gpt-4.1-nano slicing · motivates why we need KG next")

    qa_meta = _load_qa_meta()
    v = _slice_f1(REPO / "data/eval/vanilla/gpt-4.1-nano/predictions.jsonl", qa_meta)
    t = _slice_f1(REPO / "data/eval/timefilter/gpt-4.1-nano/predictions.jsonl", qa_meta)

    def avg(xs):
        return sum(xs) / len(xs) if xs else 0

    # Hop table
    hop_rows = [["Hop count", "Vanilla F1", "TimeFilter F1", "Δ", "n"]]
    for hop in [1, 2, 3]:
        key = ("hop", hop)
        vv, tt = v.get(key, []), t.get(key, [])
        if not vv:
            continue
        hop_rows.append([f"hop = {hop}", f"{avg(vv):.3f}",
                         f"{avg(tt):.3f}", f"{avg(tt)-avg(vv):+.3f}",
                         str(len(vv))])
    _add_table(s, Inches(0.4), Inches(1.4), Inches(6.0), Inches(2.0),
               hop_rows, col_widths=[Inches(1.5), Inches(1.2), Inches(1.3),
                                     Inches(1.0), Inches(1.0)], font_size=12)

    # Scope table
    scope_rows = [["Scope", "Vanilla F1", "TimeFilter F1", "Δ", "n"]]
    priority = ["inter_year", "cross_company", "intra",
                "fiscal_vs_calendar", "forward_looking"]
    for sc in priority:
        key = ("scope", sc)
        vv, tt = v.get(key, []), t.get(key, [])
        if not vv:
            continue
        scope_rows.append([sc, f"{avg(vv):.3f}", f"{avg(tt):.3f}",
                           f"{avg(tt)-avg(vv):+.3f}", str(len(vv))])
    _add_table(s, Inches(6.8), Inches(1.4), Inches(6.4), Inches(2.7),
               scope_rows, col_widths=[Inches(2.0), Inches(1.1), Inches(1.3),
                                       Inches(1.0), Inches(1.0)], font_size=12)

    # Observations
    _add_bullet_box(s, Inches(0.4), Inches(4.4), Inches(12.6), Inches(2.6), [
        ("TimeFilter shines on year-sensitive questions",
         "— inter_year +0.055, cross_company +0.051.  "
         "Exactly where a year mask should help."),
        ("No lift on intra (same-year) questions",
         "— Δ ≈ 0.  Year filter is a no-op when the question and all "
         "candidate chunks share one year."),
        ("hop = 3 barely moves",
         "— Δ = +0.007 (n=15).  Multi-hop reasoning needs entity "
         "cross-linking, which metadata alone cannot provide."),
        ("→  This motivates KG²RAG and TempoRAG-KG",
         "— graph walks add the cross-chunk relations TimeFilter "
         "is blind to."),
    ], font_size=12)


def slide_7_next(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "What's next — KG extraction and the L2 / L3 lanes")

    # KG status (left)
    lbox = s.shapes.add_textbox(Inches(0.4), Inches(1.4), Inches(6.0), Inches(0.4))
    p = lbox.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "KG extraction (running)"
    r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = NAVY

    kg_rows = [
        ["Stage", "Detail"],
        ["Corpus", "7,467 chunks from 25 × 10-K filings"],
        ["Extractor", "gpt-4.1-nano, temporal-aware prompt"],
        ["Schema", "(subject, predicate, object, valid_from, valid_to)"],
        ["Progress", "Running (live + cached), mid-sweep"],
        ["Budget", "Well under $20 cap"],
    ]
    _add_table(s, Inches(0.4), Inches(1.9), Inches(6.0), Inches(3.4),
               kg_rows, col_widths=[Inches(1.6), Inches(4.4)], font_size=12)

    # Remaining work (right)
    rbox = s.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6.4), Inches(0.4))
    p = rbox.text_frame.paragraphs[0]
    r = p.add_run(); r.text = "Remaining work"
    r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = NAVY

    _add_bullet_box(s, Inches(6.8), Inches(1.9), Inches(6.4), Inches(4.8), [
        ("1.  Finish extraction",
         "— complete full 7,467-chunk pass."),
        ("2.  Filter pass",
         "— drop noisy triples "
         "(e.g. object=\"true\", comparative artifacts)."),
        ("3.  Build KG²RAG retriever",
         "— seed chunk → graph walk → expanded context."),
        ("4.  Build TempoRAG-KG retriever",
         "— KG²RAG + hard mask on valid_from / valid_to."),
        ("5.  Run L2 + L3 sweeps",
         "— same 7 models × 129 Qs."),
        ("6.  Interaction analysis",
         "— hop × scope × condition, test whether L3 > L1 + L2."),
    ], font_size=12)

    # Footer note (no deadline emphasis)
    fbox = s.shapes.add_textbox(Inches(0.4), Inches(6.7), Inches(12.8), Inches(0.4))
    fp = fbox.text_frame.paragraphs[0]
    fp.text = ("Current signal (L1 alone) already supports RQ4 directionally; "
               "L3 will test whether the effect is multiplicative or additive.")
    fp.runs[0].font.size = Pt(12)
    fp.runs[0].font.italic = True
    fp.runs[0].font.color.rgb = MUTED


def slide_8_takeaways(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _title(s, "Takeaways",
           "From this sprint · L0 + L1 baselines on 10-K temporal QA")
    _add_bullet_box(s, Inches(0.6), Inches(1.6), Inches(12.2), Inches(5.3), [
        ("1.  TimeFilter baseline proven",
         "— +5.6% to +24.7% F1 across 7 models; "
         "lift is universal, not model-specific."),
        ("2.  Retrieval, not answer quality, is the bottleneck",
         "— F1@answered stays flat around 0.35 while coverage moves "
         "+3 to +10pp.  Any further lift must come from better retrieval."),
        ("3.  Small models benefit more than large",
         "— llama-8b +24.7% vs gpt-4o-mini +7.8%.  "
         "Directional support for the RQ4 capability-parity story."),
        ("4.  Multi-hop is next",
         "— TimeFilter cannot lift hop = 3 (+0.007).  "
         "KG²RAG and TempoRAG-KG will target exactly this gap."),
    ], font_size=14)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_1_recap(prs)
    slide_2_rqs(prs)
    slide_3_pipeline(prs)
    slide_5_results(prs)
    slide_6_ceiling(prs)
    slide_7_next(prs)
    slide_8_takeaways(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"wrote {OUT.relative_to(REPO)}  ({OUT.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
